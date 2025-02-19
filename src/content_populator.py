from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import logging
from typing import Dict, Tuple
import os

class ContentPopulator:
    """内容填充器，负责将内容填充到PPT模板中"""
    
    def __init__(self, template_path: str):
        """
        初始化内容填充器
        
        Args:
            template_path: PPT模板文件路径
        """
        self.prs = Presentation(template_path)
        
    def _get_layout_by_name(self, layout_name: str):
        """
        根据名称获取幻灯片布局
        
        Args:
            layout_name: 布局名称
            
        Returns:
            布局对象或 None
        """
        # 使用第一个可用的母版
        if not self.prs.slide_masters:
            raise ValueError("PPT模板中没有找到任何母版，请检查PPT模板是否正确")
        
        main_master = self.prs.slide_masters[0]
        logging.info(f"\n可用的母版: {main_master.name}")
        
        # 记录所有可用的布局
        logging.info("\n所有可用的布局:")
        for idx, layout in enumerate(main_master.slide_layouts):
            logging.info(f"  {idx}. {layout.name}")
            logging.info("     占位符:")
            for ph in layout.placeholders:
                ph_type = ph.placeholder_format.type
                logging.info(f"       - ID: {ph.placeholder_format.idx}, Type: {ph_type}, Name: {ph.name}")
        
        # 在母版中查找指定布局
        layout = None
        for slide_layout in main_master.slide_layouts:
            if slide_layout.name == layout_name:
                layout = slide_layout
                break
        
        if not layout:
            # 如果找不到指定布局，使用第一个可用的布局
            if main_master.slide_layouts:
                layout = main_master.slide_layouts[0]
                logging.info(f"\n找不到布局 '{layout_name}'，使用默认布局: {layout.name}")
            else:
                raise ValueError(f"在母版中没有找到任何可用的布局")
        else:
            logging.info(f"\n使用布局: {layout.name}")
        
        # 记录当前布局的占位符
        logging.info("\n当前布局的占位符:")
        for ph in layout.placeholders:
            ph_type = ph.placeholder_format.type
            logging.info(f"  - ID: {ph.placeholder_format.idx}, Type: {ph_type}, Name: {ph.name}")
        
        return layout
    
    def add_slide(self, layout_name: str, title: str = None) -> Tuple:
        """
        添加新的幻灯片
        
        Args:
            layout_name: 布局名称
            title: 幻灯片标题（可选）
            
        Returns:
            Tuple: (幻灯片对象, 布局对象)
        """
        try:
            # 获取布局
            layout = self._get_layout_by_name(layout_name)
            if not layout:
                raise ValueError(f"找不到布局: {layout_name}")
            
            # 添加新幻灯片
            slide = self.prs.slides.add_slide(layout)
            
            # 设置标题
            if title:
                self.fill_title(slide, title)
            
            return slide, layout
            
        except Exception as e:
            logging.error(f"添加幻灯片时发生错误: {str(e)}")
            raise
    
    def fill_image(self, slide, placeholder_idx: int, image_path: str):
        """
        填充图片到占位符
        
        Args:
            slide: 幻灯片对象
            placeholder_idx: 占位符索引
            image_path: 图片路径
        """
        try:
            # 检查并记录可用的占位符
            content_placeholders = []
            logging.info("\n开始检查幻灯片中的所有占位符:")
            for shape in slide.placeholders:
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                logging.info(f"发现占位符 - ID: {ph_idx}, Type: {ph_type}, Name: {shape.name}")
                if ph_type == 18:  # 只使用 PICTURE(18) 类型的占位符
                    content_placeholders.append({
                        'idx': ph_idx,
                        'type': ph_type,
                        'name': shape.name
                    })
                    logging.info(f"    -> 这是一个可用的图片占位符")
            
            logging.info(f"\n准备插入图片: {image_path}")
            print(f"当前处理的占位符索引: {placeholder_idx}")
            print(f"可用的图片占位符: {[p['idx'] for p in content_placeholders]}")
            
            # 如果没有找到合适的图片占位符，记录警告并返回
            if not content_placeholders:
                logging.warning(f"幻灯片上没有找到图片占位符(Type=18)，跳过图片插入")
                return
            
            # 根据placeholder_idx选择对应的图片占位符
            target_placeholder = None
            if placeholder_idx < len(content_placeholders):
                ph_info = content_placeholders[placeholder_idx]
                target_placeholder = slide.placeholders[ph_info['idx']]
                logging.info(f"使用第 {placeholder_idx} 个图片占位符 (ID: {ph_info['idx']}, Name: {ph_info['name']})")
            else:
                logging.warning(f"图片索引 {placeholder_idx} 超出可用占位符数量 {len(content_placeholders)}，跳过插入")
                return
            
            # 获取图片尺寸
            with Image.open(image_path) as img:
                width, height = img.size
            
            # 计算缩放比例
            target_width = target_placeholder.width
            target_height = target_placeholder.height
            scale = min(
                target_width / width,
                target_height / height
            )
            
            # 计算缩放后的尺寸
            scaled_width = int(width * scale)
            scaled_height = int(height * scale)
            
            # 计算居中位置
            # 在占位符内居中
            left = target_placeholder.left + (target_width - scaled_width) / 2
            top = target_placeholder.top + (target_height - scaled_height) / 2
            
            # 清除占位符中的任何现有内容
            if hasattr(target_placeholder, 'text'):
                target_placeholder.text = ''
            
            # 在占位符的位置插入图片形状
            pic = slide.shapes.add_picture(
                image_path,
                left,
                top,
                width=scaled_width,
                height=scaled_height
            )
            
            logging.info(f"图片插入成功: {image_path}")
            logging.info(f"图片尺寸: 原始({width}x{height}) -> 缩放后({scaled_width}x{scaled_height})")
            
        except Exception as e:
            logging.error(f"填充图片时发生错误: {str(e)}")
            raise
    
    def fill_text(self, slide, placeholder_idx: int, text_content: str):
        """
        填充文本到占位符
        
        Args:
            slide: 幻灯片对象
            placeholder_idx: 占位符索引
            text_content: 文本内容
        """
        try:
            # 检查并记录可用的占位符
            content_placeholders = []
            logging.info("\n开始检查幻灯片中的文本占位符:")
            for shape in slide.placeholders:
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                logging.info(f"发现占位符 - ID: {ph_idx}, Type: {ph_type}, Name: {shape.name}")
                if ph_type == 2:  # 只使用 BODY(2) 类型的占位符
                    content_placeholders.append({
                        'idx': ph_idx,
                        'type': ph_type,
                        'name': shape.name
                    })
                    logging.info(f"    -> 这是一个可用的文本占位符")
            
            logging.info(f"\n准备插入文本，长度: {len(text_content)} 字符")
            print(f"当前处理的占位符索引: {placeholder_idx}")
            print(f"可用的文本占位符: {[p['idx'] for p in content_placeholders]}")
            
            # 如果没有找到合适的文本占位符，记录警告并返回
            if not content_placeholders:
                logging.warning(f"幻灯片上没有找到文本占位符(Type=2)，跳过文本插入")
                return
            
            # 根据placeholder_idx选择对应的文本占位符
            target_placeholder = None
            if placeholder_idx < len(content_placeholders):
                ph_info = content_placeholders[placeholder_idx]
                target_placeholder = slide.placeholders[ph_info['idx']]
                logging.info(f"使用第 {placeholder_idx} 个文本占位符 (ID: {ph_info['idx']}, Name: {ph_info['name']})")
            else:
                logging.warning(f"文本索引 {placeholder_idx} 超出可用占位符数量 {len(content_placeholders)}，跳过插入")
                return
            
            # 插入文本内容
            if target_placeholder is not None:
                target_placeholder.text = text_content
                logging.info(f"文本插入成功")
            
        except Exception as e:
            logging.error(f"填充文本时发生错误: {str(e)}")
            raise
    
    def fill_video(self, slide, placeholder_idx: int, video_path: str):
        """
        填充视频到占位符
        
        Args:
            slide: 幻灯片对象
            placeholder_idx: 占位符索引
            video_path: 视频文件路径
        """
        try:
            # 检查并记录可用的占位符
            content_placeholders = []
            logging.info("\n开始检查幻灯片中的媒体占位符:")
            for shape in slide.placeholders:
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                logging.info(f"发现占位符 - ID: {ph_idx}, Type: {ph_type}, Name: {shape.name}")
                if ph_type == 10:  # 只使用 MEDIA_CLIP(10) 类型的占位符
                    content_placeholders.append({
                        'idx': ph_idx,
                        'type': ph_type,
                        'name': shape.name
                    })
                    logging.info(f"    -> 这是一个可用的媒体占位符")
            
            logging.info(f"\n准备插入视频: {video_path}")
            print(f"当前处理的占位符索引: {placeholder_idx}")
            print(f"可用的媒体占位符: {[p['idx'] for p in content_placeholders]}")
            
            # 如果没有找到合适的媒体占位符，记录警告并返回
            if not content_placeholders:
                logging.warning(f"幻灯片上没有找到媒体占位符(Type=10)，跳过视频插入")
                return
            
            # 根据placeholder_idx选择对应的媒体占位符
            target_placeholder = None
            if placeholder_idx < len(content_placeholders):
                ph_info = content_placeholders[placeholder_idx]
                target_placeholder = slide.placeholders[ph_info['idx']]
                logging.info(f"使用第 {placeholder_idx} 个媒体占位符 (ID: {ph_info['idx']}, Name: {ph_info['name']})")
            else:
                logging.warning(f"视频索引 {placeholder_idx} 超出可用占位符数量 {len(content_placeholders)}，跳过插入")
                return
            
            # 插入视频
            if target_placeholder is not None:
                # 清除占位符中的任何现有内容
                if hasattr(target_placeholder, 'text'):
                    target_placeholder.text = ''
                
                # 获取视频的绝对路径
                video_abs_path = os.path.abspath(video_path)
                
                # 在占位符的位置插入视频
                movie = slide.shapes.add_movie(
                    video_abs_path,
                    target_placeholder.left,
                    target_placeholder.top,
                    target_placeholder.width,
                    target_placeholder.height,
                    poster_frame_image=None,  # 可以设置视频封面图
                    mime_type='video/mp4'  # 设置适当的 MIME 类型
                )
                
                logging.info(f"视频插入成功: {video_path}")
            
        except Exception as e:
            logging.error(f"填充视频时发生错误: {str(e)}")
            raise
    
    def fill_title(self, slide, title: str):
        """
        填充幻灯片标题
        
        Args:
            slide: 幻灯片对象
            title: 标题文本
        """
        try:
            # 直接使用ID为0的占位符作为标题占位符
            try:
                print(f"\n尝试设置标题: {title}")
                print("可用的占位符:")
                for ph in slide.placeholders:
                    print(f"  - ID: {ph.placeholder_format.idx}, Type: {ph.placeholder_format.type}, Name: {ph.name}")
                
                title_placeholder = slide.placeholders[0]  # 标题占位符的ID通常是0
                print(f"找到标题占位符: ID={title_placeholder.placeholder_format.idx}, Type={title_placeholder.placeholder_format.type}, Name={title_placeholder.name}")
                
                title_placeholder.text = title
                print(f"标题设置成功: {title}")
                logging.info(f"已设置幻灯片标题: {title}")
            except (KeyError, AttributeError) as e:
                print(f"设置标题失败: {str(e)}")
                logging.warning(f"未找到标题占位符(ID=0)，无法设置标题: {str(e)}")
                
        except Exception as e:
            print(f"设置标题时发生错误: {str(e)}")
            logging.error(f"设置标题时发生错误: {str(e)}")
            raise
    
    def save(self, output_path: str):
        """
        保存PPT文件
        
        Args:
            output_path: 输出文件路径
        """
        try:
            self.prs.save(output_path)
        except Exception as e:
            logging.error(f"保存PPT时发生错误: {str(e)}")
            raise
