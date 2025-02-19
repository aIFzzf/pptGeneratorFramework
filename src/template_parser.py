from pptx import Presentation
from typing import Dict, List, Tuple
import logging

class TemplateParser:
    """PPT模板解析器，用于解析母版模板中的布局和占位符"""
    
    def __init__(self, template_path: str):
        """
        初始化模板解析器
        
        Args:
            template_path: PPT模板文件路径
        """
        self.template_path = template_path
        self.prs = None
        self.layouts_info = {}
        
    def parse(self) -> Dict:
        """
        解析PPT模板，提取布局和占位符信息
        
        Returns:
            Dict: 包含布局和占位符信息的字典
        """
        try:
            self.prs = Presentation(self.template_path)
            
            # 使用第一个可用的母版
            if not self.prs.slide_masters:
                raise ValueError("PPT模板中没有找到任何母版，请检查PPT模板是否正确")
            
            main_master = self.prs.slide_masters[0]
            logging.info(f"使用母版: {main_master.name}")
            
            # 解析该母版下的所有布局
            for layout in main_master.slide_layouts:
                placeholders = []
                for shape in layout.placeholders:
                    placeholder_info = {
                        'type': shape.placeholder_format.type,
                        'position': (shape.left, shape.top),
                        'size': (shape.width, shape.height),
                        'idx': shape.placeholder_format.idx
                    }
                    placeholders.append(placeholder_info)
                
                self.layouts_info[layout.name] = {
                    'placeholders': placeholders,
                    'placeholder_count': len(placeholders)
                }
            
            return self.layouts_info
            
        except Exception as e:
            logging.error(f"解析模板时发生错误: {str(e)}")
            raise
    
    def get_layout_by_content(self, content_type: str, count: int) -> str:
        """
        根据内容类型和数量获取合适的布局
        
        Args:
            content_type: 内容类型（image, text, video）
            count: 内容数量
            
        Returns:
            str: 布局名称
        """
        # 实现布局匹配逻辑
        pass
